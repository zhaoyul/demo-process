#!/usr/bin/env python3
"""红创科技 — 投标PPT: 嵌入仿真动画视频"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

ROOT = Path(__file__).parent
RENDERS = ROOT / "renders"
OUT = ROOT / "红创科技仿真平台_投标讲稿.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

RED = RGBColor(0xC6, 0x28, 0x28)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
DARK = RGBColor(0x1A, 0x1A, 0x2E)

def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def add_title(slide, text, y=Inches(0.2)):
    txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RED

def add_subtitle(slide, text, y=Inches(0.85)):
    txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

def add_video(slide, path, left, top, width=None, height=None):
    """Embed video from file"""
    path = str(path)
    if width is None:
        width = Inches(12)
    if height is None:
        height = Inches(5.5)
    # python-pptx video support: add as movie
    try:
        slide.shapes.add_movie(path, left, top, width, height, 
                               poster_frame_image=None, mime_type='video/mp4')
    except Exception as e:
        # Fallback: add text note
        txBox = slide.shapes.add_textbox(left, top, width, height)
        p = txBox.text_frame.paragraphs[0]
        p.text = f"[视频: {Path(path).name}]\n请在PowerPoint中手动插入此视频"
        p.font.size = Pt(16)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER

# ================================================================
# Slide 1: Title
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_bg(sl)

txBox = sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.2))
p = txBox.text_frame.paragraphs[0]
p.text = "红创科技 多物理场仿真平台"
p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = RED; p.alignment = PP_ALIGN.CENTER

txBox = sl.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(0.8))
p = txBox.text_frame.paragraphs[0]
p.text = "仿真动画演示 — 投标技术讲稿"
p.font.size = Pt(24); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

txBox = sl.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.5))
p = txBox.text_frame.paragraphs[0]
p.text = "2026年5月"
p.font.size = Pt(16); p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER

# ================================================================
# Slides 2-5: Each with one video
# ================================================================
videos = [
    # Demo 1: 悬臂梁静力学 (Hex8+B-bar+集中力)
    ("Demo 1: 悬臂梁静力学 (Hex8 + B-bar + 集中力)", "beam_loading.mp4",
     "Hex8 六面体单元 · B-bar 体积锁定处理 · 集中载荷 · 3000× 变形放大 · 逐面着色"),
    # Demo 2: 接触力学 (5s真实FEM渲染)
    ("Demo 2: 接触力学 (5s 真实 FEM 渲染)", "contact.mp4",
     "Lagrange 乘子法接触约束 · Coulomb 摩擦 · 位移 0→0.17mm · 网格逐面着色 · 真实求解器渲染"),
    # Demo 3: 双材料静电场 (含路径采样)
    ("Demo 3: 双材料静电场", "electrostatic.mp4",
     "钢(σ=10⁷ S/m) + 混凝土(σ=10⁻² S/m) · 电位 0→1V 渐进 · 3D 挤出网格", "electrostatic_path.mp4"),
    # Demo 4: 热-力-损伤三场耦合 (CDP 章节)
    ("Demo 4: 热-力-损伤三场耦合 (CDP)", "cdp_damage.mp4",
     "Concrete Damaged Plasticity · 温度-应力-损伤直接耦合 · 损伤变量演化 · 单求解器全隐式"),
]

for entry in videos:
    if len(entry) == 4:
        title, fname, desc, extra_fname = entry
        has_extra = True
    else:
        title, fname, desc = entry
        has_extra = False

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(sl)
    add_title(sl, title)
    add_subtitle(sl, desc)
    
    vpath = RENDERS / fname
    if vpath.exists():
        if has_extra:
            # Main video (left 75%)
            add_video(sl, vpath, Inches(0.3), Inches(1.3), Inches(9.5), Inches(5.8))
            # Path sampling video (right 25%, smaller)
            extra_path = RENDERS / extra_fname
            if extra_path.exists():
                txBox = sl.shapes.add_textbox(Inches(10.1), Inches(1.3), Inches(3), Inches(0.4))
                p = txBox.text_frame.paragraphs[0]
                p.text = "路径采样"
                p.font.size = Pt(12); p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
                add_video(sl, extra_path, Inches(10.1), Inches(1.8), Inches(3), Inches(5.3))
            else:
                txBox = sl.shapes.add_textbox(Inches(10.1), Inches(3), Inches(3), Inches(1))
                p = txBox.text_frame.paragraphs[0]
                p.text = f"[路径采样未找到]"
                p.font.size = Pt(14); p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER
        else:
            add_video(sl, vpath, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8))
    else:
        txBox = sl.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"[视频未找到: {fname}]"
        p.font.size = Pt(20); p.font.color.rgb = RED; p.alignment = PP_ALIGN.CENTER

# ================================================================
# Slide 7: Summary
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

txBox = sl.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(0.8))
p = txBox.text_frame.paragraphs[0]
p.text = "平台能力总结"
p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = RED; p.alignment = PP_ALIGN.CENTER

items = [
    "✓ 4 大演示场景: 静力学 · 接触 · 静电 · 热-力-损伤(CDP)",
    "✓ Hex8 + B-bar 单元 · Lagrange 乘子法接触约束",
    "✓ 温度-应力-损伤三场直接耦合 (CDP 本构)",
    "✓ 网格收敛精度 0.11% (vs Euler-Bernoulli 理论解)",
    "✓ 端到端管线: Gmsh → hongchuang-opt (MOOSE) → ParaView",
    "✓ 11 个 ExodusII (.e) 输出文件 · 完整文档体系",
]
for i, item in enumerate(items):
    txBox = sl.shapes.add_textbox(Inches(2), Inches(1.8 + i * 0.65), Inches(9), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(18); p.font.color.rgb = GOLD if i == 2 else WHITE

# Save
prs.save(str(OUT))
print(f"PPT saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
